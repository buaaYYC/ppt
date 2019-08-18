from pptx import Presentation
import ppt1
import extract_image
from text_extract import article_extract
from textrank4zh import TextRank4Keyword, TextRank4Sentence
"""
ppt 生成，本方法有局限，适用于新闻，知乎等小文章
"""
#extract text from web

def textRank_ppt(url,num_abs):
    """

    :param url:
    :param num_abs: 生成ppt张数
    :return:
    """
    title,texts = article_extract(url)
    tr4w = TextRank4Keyword()
    tr4w.analyze(text=texts ,lower=True, window=2)   # py2中text必须是utf8编码的str或者unicode对象，py3中必须是utf8编码的bytes或者str对象
    tr4s = TextRank4Sentence()
    tr4s.analyze(text=texts, lower=True, source = 'all_filters')
    print( '关键词：' )
    key_words = ""
    for item in tr4w.get_keywords(7, word_min_len=2):
        print(item.word, item.weight)
        key_words = key_words + item.word + "\n"

    #ppt generate
    prs = Presentation()
    slide1, body_shape1 = ppt1.add_slide(prs=prs,slide_title=title, style_number=0)
    slide2, body_shape2 = ppt1.add_slide(prs=prs, style_number=1,slide_title="关键词",content="")
    ppt1.add_paragraph(body_shape2,text=key_words,size=20)
    i = 0
    #图片生成，并添加到ppt中
    extract_image.pic_extract(url)
    print("句子：")
    for item in tr4s.get_key_sentences(num=(num_abs-2)*2):

        if i % 2 == 0:
            slide3, body_shape3 = ppt1.add_slide(prs=prs, style_number=1,slide_title="摘要",content="")
            try:
                ppt1.add_picture(slide2=slide3,pic_path="image1/image_"+str(i)+".jpg")
            except:
                print("no picture")
        i += 1
        # print(len(item.sentence),item.index)
        ppt1.add_paragraph(body_shape3,text=item.sentence,size=20)
    prs.save('test.pptx')
    print("ppt 已生成")
if __name__=="__main__":
    # url ="https://3w.huanqiu.com/a/c302d6/7Pd22MA8o4E?agt=20&tt_group_id=6725973518791475725"
    # textRank_ppt(url,6)
    pass