from pptx import Presentation
from pptx.util import Pt
from pptx.util import Inches
import datetime
import os
today = datetime.date.today()
prs1 = Presentation()
def add_blank(prs=prs1,style_number = 6):
    slide = prs.slides.add_slide(prs.slide_layouts[style_number])
    return slide
def add_slide(prs=prs1,style_number = 0,slide_title="Hello, World!",content="日期："+str(today)):
    """
    新建ppt
    例slide_layouts[1]为带标题和正文框的ppt，slide_layouts[6]为空白页ppt
    """
    slide = prs.slides.add_slide(prs.slide_layouts[style_number])
    # 对ppt的修改
    body_shape = slide.shapes.placeholders  # body_shape为本页ppt中所有shapes
    body_shape[0].text = slide_title  # 在第一个文本框中文字框架内添加文字\
    body_shape[1].text = content  # 在第二个文本框中文字框架内添加文字
    # prs.save('test.pptx')
    return slide,body_shape

def add_paragraph(body_shape,text= 'add_paragraph',size = 15):
    # 在文本框中添加新段落
    new_paragraph = body_shape[1].text_frame.add_paragraph()  # 在第二个shape中的文本框架中添加新段落
    new_paragraph.text = text  # 新段落中文字
    new_paragraph.font.bold = False  # 文字加粗
    new_paragraph.font.italic = False  # 文字斜体
    new_paragraph.font.size = Pt(size)  # 文字大小
    # new_paragraph.font.underline = True  # 文字下划线
    # new_paragraph.level = 1  # 新段落的级别

def add_textbox(slide1,text='this is a new textbox',box = [5,5]):
    # 添加新文本框
    left = top = Inches(box[0])
    width = height = Inches(box[1]) # 预设位置及大小
    textbox = slide1.shapes.add_textbox(left, top, width, height)  # left，top为相对位置，width，height为文本框大小
    textbox.text = text  # 文本框中文字
    return textbox
def add_picture(slide2,pic_path='Tulips.jpg'):
    #添加图片
    img_path = pic_path  # 文件路径
    left, top, width, height = Inches(4), Inches(4.5), Inches(4), Inches(2.5)  # 预设位置及大小
    pic = slide2.shapes.add_picture(img_path, left, top, width, height)  # 在指定位置按预设值添加图片
if __name__=="__main__":
    # prs = Presentation()
    # slide2,body_shape2 = add_slide(prs=prs,style_number=1)
    # add_paragraph(body_shape2)
    # # add_picture(slide2,pic_path=)
    # prs.save('test1.pptx')
    pass




